import streamlit as st
import openai
import tempfile
import os
import requests
import torch
from docx import Document
from pptx import Presentation
from pptx.util import Inches
from PIL import Image
from io import BytesIO
from transformers import pipeline

def extract_text_from_docx(docx_path):
    """Extract text from a Word (.docx) file."""
    doc = Document(docx_path)
    text = [para.text for para in doc.paragraphs if para.text.strip()]
    return text

def summarize_text(text, api_key):
    """Summarize text using OpenAI API."""
    if not api_key:
        return text  # If no API key, return original text
    try:
        openai.api_key = api_key
        response = openai.ChatCompletion.create(
            model="gpt-4",
            messages=[{"role": "user", "content": f"Summarize this: {text}"}]
        )
        return response['choices'][0]['message']['content'].strip()
    except Exception as e:
        st.error(f"OpenAI API Error: {str(e)}")
        return text

def generate_image(prompt, api_key):
    """Generate an image using OpenAI DALL·E API."""
    try:
        openai.api_key = api_key
        response = openai.Image.create(
            prompt=prompt,
            n=1,
            size="512x512"
        )
        image_url = response['data'][0]['url']
        image_response = requests.get(image_url)
        return Image.open(BytesIO(image_response.content))
    except Exception as e:
        st.error(f"Image Generation Error: {str(e)}")
        return None

def transcribe_audio(audio_path):
    """Convert speech to text using Whisper AI."""
    try:
        transcriber = pipeline("automatic-speech-recognition", model="openai/whisper-small")
        transcription = transcriber(audio_path)["text"]
        return transcription
    except Exception as e:
        st.error(f"Speech-to-Text Error: {str(e)}")
        return ""

def create_ppt_from_text(text_list, api_key, output_path):
    """Generate a PowerPoint file with AI-enhanced slides."""
    prs = Presentation()
    for text in text_list:
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title & Content layout
        title = slide.shapes.title
        title.text = "AI-Generated Slide"
        
        img = generate_image(text, api_key)
        if img:
            img_path = output_path.replace(".pptx", ".png")
            img.save(img_path)
            slide.shapes.add_picture(img_path, Inches(1), Inches(1), width=Inches(4), height=Inches(3))
        
        content = slide.placeholders[0]
        content.text = text
    prs.save(output_path)

def main():
    st.title("AI-Powered Word to PowerPoint Converter")
    st.write("Upload a Word document and generate a PowerPoint with AI-enhanced content and images.")
    
    user_api_key = st.text_input("Enter your OpenAI API Key (Required for AI features)", type="password")
    uploaded_file = st.file_uploader("Upload a .docx file", type=["docx"])
    
    if uploaded_file:
        with tempfile.NamedTemporaryFile(delete=False, suffix=".docx") as temp_docx:
            temp_docx.write(uploaded_file.read())
            temp_docx_path = temp_docx.name
        
        extracted_text = extract_text_from_docx(temp_docx_path)
        summarized_text = [summarize_text(text, user_api_key) for text in extracted_text]
        
        pptx_path = temp_docx_path.replace(".docx", ".pptx")
        create_ppt_from_text(summarized_text, user_api_key, pptx_path)
        
        st.success("PowerPoint created successfully!")
        with open(pptx_path, "rb") as pptx_file:
            st.download_button(
                label="Download PowerPoint",
                data=pptx_file,
                file_name="presentation.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )
        
        os.remove(temp_docx_path)
        os.remove(pptx_path)

if __name__ == "__main__":
    main()
